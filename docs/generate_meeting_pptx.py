# -*- coding: utf-8 -*-
"""ヴィレッジHP 打合せ資料 PPTX 生成スクリプト。
再生成: python docs/generate_meeting_pptx.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- ブランドカラー ----
INK      = RGBColor(0x16, 0x23, 0x2B)
INK_SOFT = RGBColor(0x47, 0x58, 0x5F)
BLUE     = RGBColor(0x1C, 0x6D, 0x8F)
BLUE_DEEP= RGBColor(0x11, 0x4A, 0x63)
GREEN    = RGBColor(0x2C, 0x8F, 0x66)
GROUND   = RGBColor(0xF4, 0xF6, 0xF7)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
AMBER    = RGBColor(0xB7, 0x79, 0x1F)
LINE     = RGBColor(0xDD, 0xE4, 0xE7)
FONT = "Yu Gothic"
FONT_B = "Yu Gothic"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def set_font(run, size, color, bold=False, name=FONT):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = name
    # 日本語フォントを明示
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', name)


def add_rect(slide, x, y, w, h, color, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             space_after=6, line_spacing=1.15):
    """runs: list of paragraphs, each a list of (text,size,color,bold) tuples."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (text, size, color, bold) in para:
            r = p.add_run(); r.text = text
            set_font(r, size, color, bold)
    return tb


def section_header(slide, no, title):
    add_rect(slide, Inches(0), Inches(0), SW, Inches(0.12), BLUE)
    # 番号バッジ
    badge = add_rect(slide, Inches(0.7), Inches(0.55), Inches(0.55), Inches(0.55), BLUE)
    tf = badge.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(no); set_font(r, 22, WHITE, True)
    badge.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    add_text(slide, Inches(1.5), Inches(0.5), Inches(11), Inches(0.7),
             [[(title, 26, INK, True)]], anchor=MSO_ANCHOR.MIDDLE)
    add_rect(slide, Inches(0.7), Inches(1.32), Inches(11.93), Emu(9525), LINE)


# ============ Slide 1: 表紙 ============
s = prs.slides.add_slide(BLANK)
add_rect(s, Inches(0), Inches(0), SW, SH, BLUE_DEEP)
add_rect(s, Inches(0), Inches(6.9), SW, Inches(0.6), GREEN)
# Vマーク（三角2枚）
tri1 = s.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(0.75), Inches(0.7), Inches(0.8), Inches(0.8))
tri1.rotation = 180; tri1.fill.solid(); tri1.fill.fore_color.rgb = RGBColor(0x5B,0xB2,0xD1); tri1.line.fill.background(); tri1.shadow.inherit=False
tri2 = s.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(1.0), Inches(0.95), Inches(0.32), Inches(0.5))
tri2.rotation = 180; tri2.fill.solid(); tri2.fill.fore_color.rgb = GREEN; tri2.line.fill.background(); tri2.shadow.inherit=False

add_text(s, Inches(0.8), Inches(2.5), Inches(11.7), Inches(0.5),
         [[("株式会社ヴィレッジ　ホームページ制作", 18, RGBColor(0x9E,0xC7,0xD8), True)]])
add_text(s, Inches(0.8), Inches(3.05), Inches(11.7), Inches(1.6),
         [[("打合せ資料", 54, WHITE, True)],
          [("再来週（7/20 週）", 30, RGBColor(0xCF,0xE6,0xEF), False)]], line_spacing=1.1, space_after=10)
add_text(s, Inches(0.8), Inches(5.6), Inches(11.7), Inches(0.8),
         [[("作成日 2026-07-06　|　土曜ヒアリング反映版　|　目的：双方の分担確認・決定事項の共有",
            14, RGBColor(0xAF,0xCE,0xDB), False)]])

# ============ Slide 2: サイトの意図・目的 ============
s = prs.slides.add_slide(BLANK)
add_rect(s, Inches(0), Inches(0), SW, SH, GROUND)
section_header(s, 1, "サイトの意図・目的")
points = [
    ("設立2024年の新しい会社", "だからこそ「信頼感・誠実・清潔感」を前面に出し、発注前の不安を取り除くことが最大の狙い。"),
    ("解体から活用までワンストップ", "（解体・土地開発・リフォーム・リノベ）を打ち出し、単なる解体屋ではない付加価値を伝える。"),
    ("有資格者による品質管理＋近隣対策", "（クレームの少なさ）を強みの軸に、地域密着・スピード対応をアピール。"),
    ("電話・問い合わせフォームでの反響獲得", "がゴール。問い合わせ導線を明確にする。"),
]
y = 1.75
for head, body in points:
    card = add_rect(s, Inches(0.7), Inches(y), Inches(11.93), Inches(1.05), WHITE)
    add_rect(s, Inches(0.7), Inches(y), Inches(0.1), Inches(1.05), GREEN)
    add_text(s, Inches(1.05), Inches(y+0.1), Inches(11.3), Inches(0.85),
             [[(head, 17, BLUE, True), ("　"+body, 15, INK_SOFT, False)]],
             anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.2)
    y += 1.22

# ============ Slide 3: 現サンプルHPの状況 ============
s = prs.slides.add_slide(BLANK)
add_rect(s, Inches(0), Inches(0), SW, SH, GROUND)
section_header(s, 2, "現サンプルHPの状況")
rows = [
    (GREEN, "トップ1枚構成のドラフト完成。", "「案1（洗練ブランド）」に全面刷新してmainへ統合。"),
    (GREEN, "土曜の決定事項は主要部分を反映済み。", "キャッチコピー変更／不動産売買を除外し5事業へ／許可番号・設立2024年を追記／近隣対策を「MEASURES」新設／Googleマップ埋め込み。"),
    (AMBER, "写真・ロゴは全てサンプル画像のまま。", "実データ待ち。"),
    (AMBER, "本番URLは未公開（localhostのみ）。", "確認用URLの用意が必要。"),
]
y = 1.75
for dot, head, body in rows:
    card = add_rect(s, Inches(0.7), Inches(y), Inches(11.93), Inches(1.12), WHITE)
    d = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), Inches(y+0.28), Inches(0.22), Inches(0.22))
    d.fill.solid(); d.fill.fore_color.rgb = dot; d.line.fill.background(); d.shadow.inherit=False
    add_text(s, Inches(1.45), Inches(y+0.13), Inches(10.9), Inches(0.9),
             [[(head, 16, INK, True)], [(body, 14, INK_SOFT, False)]],
             anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15, space_after=2)
    y += 1.28

# ============ Slide 4: 決定事項テーブル ============
s = prs.slides.add_slide(BLANK)
add_rect(s, Inches(0), Inches(0), SW, SH, GROUND)
section_header(s, 3, "土曜ヒアリング 決定事項")
data = [
    ("キャッチコピー", "解体工事で未来のかけ橋を創造する", "済"),
    ("事業内容", "不動産売買は今後の新規事業のため非掲載 → 5事業", "済"),
    ("補助金申請代行", "追加事業・対応市区町村は未記入", "保留"),
    ("許可番号", "解体工事業登録 第1456号／石綿調査者 第21071177号", "済"),
    ("選ばれる理由①②③", "そのまま採用でOK", "済"),
    ("他社に負けない点", "近隣対策（着手前・完了後のあいさつ）＝クレームの少なさ", "済"),
    ("営業時間", "8:00〜18:00／日祝休み（実際の時間は要確認）", "済"),
    ("設立", "2024年", "済"),
    ("サイトの色・雰囲気", "青×緑を維持／信頼感・誠実・清潔感", "済"),
    ("参考サイト", "㈱ACTIVE・㈱サンライズ・㈱クリーンアイランド", "研究"),
    ("問い合わせフォーム／地図", "フォーム設置（未実装）／Googleマップ掲載（済）", "一部"),
    ("ドメイン", ".jp 希望", "取得へ"),
]
tx, ty, tw = Inches(0.7), Inches(1.65), Inches(11.93)
rows_n = len(data) + 1
tbl_shape = s.shapes.add_table(rows_n, 3, tx, ty, tw, Inches(5.3))
table = tbl_shape.table
table.columns[0].width = Inches(3.1)
table.columns[1].width = Inches(7.3)
table.columns[2].width = Inches(1.53)
# ヘッダー
hdr = ["項目", "決定内容", "反映"]
for c, txt in enumerate(hdr):
    cell = table.cell(0, c)
    cell.fill.solid(); cell.fill.fore_color.rgb = BLUE_DEEP
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_top = Pt(2); cell.margin_bottom = Pt(2)
    p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT if c<2 else PP_ALIGN.CENTER
    r = p.add_run(); r.text = txt; set_font(r, 13, WHITE, True)
# 本文
badge_color = {"済":GREEN, "保留":AMBER, "研究":BLUE, "取得へ":BLUE, "一部":AMBER}
for i, (a, b, c) in enumerate(data, start=1):
    fill = WHITE if i % 2 else RGBColor(0xEC,0xF1,0xF3)
    for col, txt, sz, col_color, bold, al in [
        (0, a, 12, INK, True, PP_ALIGN.LEFT),
        (1, b, 11.5, INK_SOFT, False, PP_ALIGN.LEFT),
        (2, c, 11, badge_color.get(c, INK), True, PP_ALIGN.CENTER)]:
        cell = table.cell(i, col)
        cell.fill.solid(); cell.fill.fore_color.rgb = fill
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Pt(6); cell.margin_top = Pt(1); cell.margin_bottom = Pt(1)
        p = cell.text_frame.paragraphs[0]; p.alignment = al
        r = p.add_run(); r.text = txt; set_font(r, sz, col_color, bold)

# ============ Slide 5: 制作側で進めること ============
s = prs.slides.add_slide(BLANK)
add_rect(s, Inches(0), Inches(0), SW, SH, GROUND)
section_header(s, 4, "お互いの分担（〜打ち合わせ）")
add_rect(s, Inches(0.7), Inches(1.6), Inches(3.0), Inches(0.5), BLUE)
add_text(s, Inches(0.7), Inches(1.6), Inches(3.0), Inches(0.5),
         [[("制作側で進めること", 16, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
me = [
    ("ドメイン取得", "（.jp／village-〜 の空き確認 → 取得手続き）"),
    ("独自ドメインのメール作成", "（village2024@outlook.jp → 独自ドメインへ）"),
    ("確認用サイトを公開（Vercel等）", "し、閲覧URLを用意"),
    ("問い合わせフォーム", "の実装"),
    ("参考3サイトを研究", "しデザインを微調整"),
    ("提供写真・正式ロゴ・実営業時間を反映", "（受領後）"),
    ("確認事項を整理して持参", "（事業項目／補助金申請代行の要否）"),
]
y = 2.35
for head, body in me:
    dash = add_rect(s, Inches(0.75), Inches(y+0.14), Inches(0.16), Inches(0.16), BLUE)
    add_text(s, Inches(1.1), Inches(y), Inches(11.5), Inches(0.5),
             [[(head, 14.5, INK, True), (body, 13, INK_SOFT, False)]], line_spacing=1.1)
    y += 0.62

# ============ Slide 6: お客様にご用意いただきたいこと ============
s = prs.slides.add_slide(BLANK)
add_rect(s, Inches(0), Inches(0), SW, SH, GROUND)
section_header(s, 4, "お客様にご用意いただきたいこと")
add_rect(s, Inches(0.7), Inches(1.55), Inches(5.6), Inches(0.45), GREEN)
add_text(s, Inches(0.7), Inches(1.55), Inches(5.6), Inches(0.45),
         [[("打ち合わせまでに最低限そろえてほしい情報", 14, WHITE, True)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
must = [
    ("写真データ一式", "（下記の内訳・スマホ撮影でOK）"),
    ("会社ロゴの正式データ", ""),
    ("正確な営業時間・定休日", "（8:00〜18:00／日祝休みで合っているか）"),
    ("掲載する事業内容の最終確定", "（不動産売買は除外／補助金申請代行の要否）"),
    ("会社概要の確定情報", "（資本金・従業員数・その他許認可の有無）"),
    ("表示する連絡先の最終確認", "（電話・メールを載せてよいか）"),
    ("希望ドメインの文字列", "（village-〇〇 の〇〇。末尾 .jp）"),
]
y = 2.2
for head, body in must:
    chk = add_rect(s, Inches(0.75), Inches(y+0.06), Inches(0.2), Inches(0.2), WHITE, line=GREEN)
    add_text(s, Inches(1.12), Inches(y), Inches(11.4), Inches(0.4),
             [[(head, 14, INK, True), (body, 12.5, INK_SOFT, False)]], line_spacing=1.05)
    y += 0.5
# 写真内訳
add_rect(s, Inches(0.7), Inches(y+0.05), Inches(11.93), Inches(0.7), RGBColor(0xE5,0xF2,0xEC))
add_text(s, Inches(0.95), Inches(y+0.1), Inches(11.5), Inches(0.6),
         [[("写真の内訳：", 12, GREEN, True),
           ("代表者 / 解体工事（前・中・後）/ 重機・車両 / 不動産・土地 / リフォーム・リノベ事例(before-after) / スタッフ・作業風景・会社外観",
            12, INK_SOFT, False)]], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15)
add_text(s, Inches(0.7), Inches(y+0.95), Inches(11.93), Inches(0.4),
         [[("あると助かるもの（任意）：参考サイトで「特に良い」と感じたポイント", 12.5, INK_SOFT, False)]])

out = "docs/ヴィレッジHP_打合せ資料.pptx"
prs.save(out)
print("saved:", out, "/ slides:", len(prs.slides._sldIdLst))
